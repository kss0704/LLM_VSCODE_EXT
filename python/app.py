import streamlit as st
import requests
import json
import time
from typing import Dict, List, Optional
import re
from datetime import datetime
import subprocess
import tempfile
import os
import sys
import zipfile
import git
from pathlib import Path
import hashlib
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import pickle
import shutil
import PyPDF2
import docx
from docx import Document
import io
import sqlite3
from concurrent.futures import ThreadPoolExecutor, as_completed
import tiktoken
import pptx
from pptx import Presentation
import openpyxl
from openpyxl import load_workbook
import csv
import xml.etree.ElementTree as ET

# Page configuration
st.set_page_config(
    page_title="LLM Code Assistant",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
<style>
.main-header {
    text-align: center;
    padding: 1rem 0;
    background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
    color: white;
    border-radius: 10px;
    margin-bottom: 2rem;
}
.chat-message {
    padding: 1rem;
    margin: 0.5rem 0;
    border-radius: 10px;
    border-left: 4px solid #667eea;
    background-color: #f8f9fa;
}
.user-message {
    border-left-color: #28a745;
    background-color: #e8f5e9;
}
.assistant-message {
    border-left-color: #667eea;
    background-color: #f0f2ff;
}
.code-block {
    background-color: #1e1e1e;
    color: #d4d4d4;
    padding: 1rem;
    border-radius: 8px;
    margin: 0.5rem 0;
    overflow-x: auto;
}
.language-badge {
    display: inline-block;
    background-color: #667eea;
    color: white;
    padding: 0.2rem 0.5rem;
    border-radius: 15px;
    font-size: 0.8rem;
    margin-bottom: 0.5rem;
}
.file-tree {
    background-color: #f8f9fa;
    padding: 1rem;
    border-radius: 8px;
    margin: 0.5rem 0;
    border-left: 4px solid #17a2b8;
}
.codebase-stats {
    background-color: #e9ecef;
    padding: 1rem;
    border-radius: 8px;
    margin: 0.5rem 0;
}
.token-warning {
    background-color: #fff3cd;
    border: 1px solid #ffeaa7;
    color: #856404;
    padding: 0.75rem;
    border-radius: 0.375rem;
    margin: 0.5rem 0;
}
.token-danger {
    background-color: #f8d7da;
    border: 1px solid #f5c6cb;
    color: #721c24;
    padding: 0.75rem;
    border-radius: 0.375rem;
    margin: 0.5rem 0;
}
.core-capabilities {
    padding: 0.5rem 0;
    margin: 1rem 0;
    border-left: 3px solid #667eea;
}
.core-capabilities h4 {
    color: #white;
    margin-bottom: 0.5rem;
    font-weight: 600;
}
.capability-item {
    display: flex;
    align-items: center;
    margin: 0.3rem 0;
    font-size: 0.85rem;
    color: #white;
    padding: 0.2rem 0;
}
.capability-icon {
    margin-right: 0.5rem;
    font-size: 1rem;
}
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'messages' not in st.session_state:
    st.session_state.messages = []
if 'api_key' not in st.session_state:
    st.session_state.api_key = ""
if 'model' not in st.session_state:
    st.session_state.model = "llama-3.1-8b-instant"
if 'temperature' not in st.session_state:
    st.session_state.temperature = 0.1
if 'max_tokens' not in st.session_state:
    st.session_state.max_tokens = 4000
if 'codebase_data' not in st.session_state:
    st.session_state.codebase_data = {}
if 'embeddings_model' not in st.session_state:
    st.session_state.embeddings_model = None
if 'code_embeddings' not in st.session_state:
    st.session_state.code_embeddings = {}
if 'codebase_handler' not in st.session_state:
    st.session_state.codebase_handler = None
if 'max_context_tokens' not in st.session_state:
    st.session_state.max_context_tokens = 120000
if 'chunk_strategy' not in st.session_state:
    st.session_state.chunk_strategy = "smart"
if 'max_file_size' not in st.session_state:
    st.session_state.max_file_size = 10 * 1024 * 1024  # 10MB
if 'file_processing_error' not in st.session_state:
    st.session_state.file_processing_error = None

# Supported programming languages
SUPPORTED_LANGUAGES = {
    'python': {'extension': '.py', 'runner': 'python'},
    'javascript': {'extension': '.js', 'runner': 'node'},
    'java': {'extension': '.java', 'runner': 'java'},
    'cpp': {'extension': '.cpp', 'runner': 'g++'},
    'c': {'extension': '.c', 'runner': 'gcc'},
    'go': {'extension': '.go', 'runner': 'go run'},
    'rust': {'extension': '.rs', 'runner': 'rustc'},
    'php': {'extension': '.php', 'runner': 'php'},
    'ruby': {'extension': '.rb', 'runner': 'ruby'},
    'swift': {'extension': '.swift', 'runner': 'swift'},
    'kotlin': {'extension': '.kt', 'runner': 'kotlinc'},
    'typescript': {'extension': '.ts', 'runner': 'ts-node'},
    'bash': {'extension': '.sh', 'runner': 'bash'},
    'powershell': {'extension': '.ps1', 'runner': 'powershell'},
    'sql': {'extension': '.sql', 'runner': 'sqlite3'},
    'html': {'extension': '.html', 'runner': 'browser'},
    'css': {'extension': '.css', 'runner': 'browser'},
    'r': {'extension': '.r', 'runner': 'Rscript'},
    'matlab': {'extension': '.m', 'runner': 'octave'}
}

# Extended file extensions to process (including all document types)
ALL_EXTENSIONS = {
    # Code files
    '.py', '.js', '.java', '.cpp', '.c', '.go', '.rs', '.php', '.rb', 
    '.swift', '.kt', '.ts', '.sh', '.ps1', '.sql', '.html', '.css', 
    '.r', '.m', '.jsx', '.tsx', '.vue', '.svelte', '.scala', '.cs', 
    '.vb', '.pl', '.lua', '.dart', '.elm', '.clj', '.hs', '.ml', 
    '.f90', '.pas', '.asm', '.bat', '.yml', '.yaml', '.json', '.xml',
    '.md', '.txt', '.cfg', '.ini', '.toml',
    # Document files
    '.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls',
    '.odt', '.ods', '.odp', '.rtf', '.tex', '.csv',
    # Image files (for OCR potential)
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',
    # Archive files
    '.zip', '.tar', '.gz', '.rar', '.7z'
}

# Groq models with context windows
GROQ_MODELS = {
    "llama-3.1-8b-instant": {"name": "Llama 3.1 8B (Fastest)", "context": 128000},
    "llama-3.1-70b-versatile": {"name": "Llama 3.1 70B (Most Capable)", "context": 128000},
    "llama-3.2-1b-preview": {"name": "Llama 3.2 1B (Preview)", "context": 128000},
    "llama-3.2-3b-preview": {"name": "Llama 3.2 3B (Preview)", "context": 128000},
    "mixtral-8x7b-32768": {"name": "Mixtral 8x7B", "context": 32000},
    "gemma2-9b-it": {"name": "Gemma 2 9B", "context": 8000},
    "gemma-7b-it": {"name": "Gemma 7B", "context": 8000},
    "llama3-8b-8192": {"name": "Llama 3 8B (Legacy)", "context": 8000},
    "llama3-70b-8192": {"name": "Llama 3 70B (Legacy)", "context": 8000}
}

# File priority for intelligent selection
FILE_PRIORITY = {
    # Code files (highest priority)
    '.py': 10, '.js': 9, '.java': 8, '.cpp': 7, '.c': 7, '.go': 6,
    '.rs': 6, '.ts': 9, '.jsx': 8, '.tsx': 8, '.vue': 7, '.svelte': 7,
    '.html': 5, '.css': 4, '.sql': 6, '.sh': 5, '.yml': 4, '.yaml': 4,
    '.json': 3, '.xml': 3, '.md': 4, '.txt': 3, '.cfg': 2, '.ini': 2,
    # Document files (medium priority)
    '.pdf': 6, '.docx': 5, '.doc': 5, '.pptx': 4, '.ppt': 4,
    '.xlsx': 5, '.xls': 5, '.csv': 7, '.odt': 4, '.rtf': 3,
    # Archive files (low priority)
    '.zip': 2, '.tar': 2, '.gz': 2,
    # Image files (lowest priority)
    '.png': 1, '.jpg': 1, '.jpeg': 1, '.gif': 1, '.bmp': 1
}

class EnhancedCodebaseHandler:
    """Enhanced handler for large codebases with support for all file types."""
    
    def __init__(self, db_path: str = None):
        self.db_path = db_path or os.path.join(tempfile.gettempdir(), "enhanced_codebase_cache.db")
        self.embeddings_model = None
        self.encoding = None
        self.init_database()
        self.init_tokenizer()
        
    def init_tokenizer(self):
        """Initialize tokenizer for accurate token counting."""
        try:
            self.encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")
        except:
            self.encoding = None
    
    def count_tokens(self, text: str) -> int:
        """Count tokens accurately."""
        if self.encoding:
            return len(self.encoding.encode(str(text)))
        else:
            return len(str(text)) // 4
    
    def init_database(self):
        """Initialize SQLite database for codebase storage."""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS files (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_path TEXT UNIQUE,
                file_hash TEXT,
                extension TEXT,
                file_type TEXT,
                size INTEGER,
                lines INTEGER,
                tokens INTEGER,
                priority INTEGER DEFAULT 1,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                last_modified TIMESTAMP,
                is_processed BOOLEAN DEFAULT FALSE,
                content_preview TEXT
            )
        """)
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS file_chunks (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_id INTEGER,
                chunk_index INTEGER,
                content TEXT,
                tokens INTEGER,
                start_line INTEGER,
                end_line INTEGER,
                chunk_type TEXT DEFAULT 'general',
                chunk_hash TEXT,
                relevance_score REAL DEFAULT 0.0,
                FOREIGN KEY (file_id) REFERENCES files (id)
            )
        """)
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS embeddings (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                chunk_id INTEGER,
                embedding BLOB,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (chunk_id) REFERENCES file_chunks (id)
            )
        """)
        
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_file_path ON files(file_path)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_file_priority ON files(priority DESC, tokens ASC)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_chunk_tokens ON file_chunks(tokens)")
        cursor.execute("CREATE INDEX IF NOT EXISTS idx_chunk_relevance ON file_chunks(relevance_score DESC)")
        
        conn.commit()
        conn.close()
    
    def extract_text_from_file(self, file_path: str) -> str:
        """Extract text from various file formats with enhanced error handling."""
        extension = Path(file_path).suffix.lower()
        
        try:
            # Text-based files
            if extension in ['.txt', '.md', '.py', '.js', '.java', '.cpp', '.c', '.go', 
                           '.rs', '.php', '.rb', '.html', '.css', '.json', '.xml', '.yml', '.yaml']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            # PDF files
            elif extension == '.pdf':
                return self.extract_pdf_text(file_path)
            
            # Word documents
            elif extension in ['.docx', '.doc']:
                return self.extract_docx_text(file_path)
            
            # PowerPoint presentations
            elif extension in ['.pptx', '.ppt']:
                return self.extract_pptx_text(file_path)
            
            # Excel files
            elif extension in ['.xlsx', '.xls']:
                return self.extract_excel_text(file_path)
            
            # CSV files
            elif extension == '.csv':
                return self.extract_csv_text(file_path)
            
            # RTF files (basic text extraction)
            elif extension == '.rtf':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            # Try as text file for any other extensions
            else:
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read()
                except:
                    return f"[Binary file: {os.path.basename(file_path)}]"
                    
        except Exception as e:
            return f"[Error reading file {os.path.basename(file_path)}: {str(e)}]"
    
    def extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF files with error handling."""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except Exception as e:
            return f"[PDF Error: {str(e)}]"
    
    def extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX files with error handling."""
        try:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text.append(cell.text)
            
            return "\n".join(text)
        except Exception as e:
            return f"[DOCX Error: {str(e)}]"
    
    def extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX files with error handling."""
        try:
            presentation = Presentation(file_path)
            text = []
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            
            return "\n".join(text)
        except Exception as e:
            return f"[PPTX Error: {str(e)}]"
    
    def extract_excel_text(self, file_path: str) -> str:
        """Extract text from Excel files with error handling."""
        try:
            workbook = load_workbook(file_path, read_only=True)
            text = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    text.append("\t".join(row_text))
            
            return "\n".join(text)
        except Exception as e:
            return f"[Excel Error: {str(e)}]"
    
    def extract_csv_text(self, file_path: str) -> str:
        """Extract text from CSV files with error handling."""
        try:
            text = []
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                csv_reader = csv.reader(file)
                for row in csv_reader:
                    text.append("\t".join(row))
            return "\n".join(text)
        except Exception as e:
            return f"[CSV Error: {str(e)}]"
    
    def smart_chunk_content(self, content: str, file_ext: str, max_tokens: int = 2000) -> List[Dict]:
        """Smart chunking that respects file structure and semantic boundaries."""
        lines = content.split('\n')
        chunks = []
        current_chunk = []
        current_tokens = 0
        start_line = 0
        chunk_type = 'general'
        
        # Language-specific patterns for semantic boundaries
        boundary_patterns = {
            '.py': [r'^\s*(def |class |import |from )', r'^\s*#.*', r'^\s*""".*"""'],
            '.js': [r'^\s*(function |class |const |let |var )', r'^\s*/\*.*\*/', r'^\s*//.*'],
            '.java': [r'^\s*(public |private |protected |class |interface )', r'^\s*/\*.*\*/', r'^\s*//.*'],
            '.cpp': [r'^\s*(class |struct |namespace |#include)', r'^\s*/\*.*\*/', r'^\s*//.*'],
            '.go': [r'^\s*(func |type |package |import)', r'^\s*/\*.*\*/', r'^\s*//.*'],
            '.md': [r'^#+ ', r'^-{3,}', r'^={3,}'],
            '.txt': [r'^[A-Z][^a-z]*$', r'^\d+\.', r'^-'],
        }
        
        patterns = boundary_patterns.get(file_ext, [r'^\s*$'])
        
        for i, line in enumerate(lines):
            line_tokens = self.count_tokens(line)
            
            # Check for semantic boundaries
            is_boundary = any(re.match(pattern, line) for pattern in patterns)
            
            # Determine chunk type
            if any(keyword in line.lower() for keyword in ['class', 'function', 'def']):
                chunk_type = 'class_function'
            elif any(keyword in line.lower() for keyword in ['import', 'include', 'require']):
                chunk_type = 'imports'
            elif line.strip().startswith(('#', '//', '/*')):
                chunk_type = 'comments'
            elif file_ext in ['.md', '.txt']:
                if line.startswith('#'):
                    chunk_type = 'heading'
                elif line.strip() == '':
                    chunk_type = 'paragraph_break'
            
            # Start new chunk at boundaries or when size limit reached
            if ((is_boundary and current_tokens > max_tokens * 0.5) or 
                (current_tokens + line_tokens > max_tokens)) and current_chunk:
                
                chunk_content = '\n'.join(current_chunk)
                chunks.append({
                    'content': chunk_content,
                    'tokens': current_tokens,
                    'start_line': start_line,
                    'end_line': start_line + len(current_chunk) - 1,
                    'chunk_type': chunk_type
                })
                
                current_chunk = [line]
                current_tokens = line_tokens
                start_line = i
            else:
                current_chunk.append(line)
                current_tokens += line_tokens
        
        # Add final chunk
        if current_chunk:
            chunk_content = '\n'.join(current_chunk)
            chunks.append({
                'content': chunk_content,
                'tokens': current_tokens,
                'start_line': start_line,
                'end_line': start_line + len(current_chunk) - 1,
                'chunk_type': chunk_type
            })
        
        return chunks
    
    def process_file(self, file_path: str) -> bool:
        """Process a single file with enhanced error handling."""
        try:
            # Check file size
            file_size = os.path.getsize(file_path)
            if file_size > st.session_state.max_file_size:
                st.warning(f"Skipping large file: {os.path.basename(file_path)} ({file_size / 1024 / 1024:.1f}MB)")
                return False
            
            # Extract content based on file type
            content = self.extract_text_from_file(file_path)
            
            # Skip if content is too large
            total_tokens = self.count_tokens(content)
            if total_tokens > 100000:  # Increased limit for documents
                st.warning(f"Skipping file with too many tokens: {os.path.basename(file_path)} ({total_tokens:,} tokens)")
                return False
            
            # Get file metadata
            file_hash = hashlib.md5(content.encode()).hexdigest()
            extension = Path(file_path).suffix.lower()
            
            # Determine file type
            if extension in ['.py', '.js', '.java', '.cpp', '.c', '.go', '.rs']:
                file_type = 'code'
            elif extension in ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls']:
                file_type = 'document'
            elif extension in ['.md', '.txt']:
                file_type = 'text'
            elif extension in ['.json', '.xml', '.yml', '.yaml']:
                file_type = 'config'
            else:
                file_type = 'other'
            
            lines = len(content.split('\n'))
            priority = FILE_PRIORITY.get(extension, 1)
            content_preview = content[:500] + "..." if len(content) > 500 else content
            
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            # Insert/update file record
            cursor.execute("""
                INSERT OR REPLACE INTO files 
                (file_path, file_hash, extension, file_type, size, lines, tokens, priority, 
                 last_modified, is_processed, content_preview)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, datetime('now'), TRUE, ?)
            """, (file_path, file_hash, extension, file_type, file_size, lines, 
                  total_tokens, priority, content_preview))
            
            file_id = cursor.lastrowid
            
            # Delete old chunks
            cursor.execute("DELETE FROM file_chunks WHERE file_id = ?", (file_id,))
            
            # Create smart chunks
            chunks = self.smart_chunk_content(content, extension)
            
            for chunk_idx, chunk in enumerate(chunks):
                chunk_hash = hashlib.md5(chunk['content'].encode()).hexdigest()
                cursor.execute("""
                    INSERT INTO file_chunks 
                    (file_id, chunk_index, content, tokens, start_line, end_line, chunk_type, chunk_hash)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                """, (file_id, chunk_idx, chunk['content'], chunk['tokens'],
                      chunk['start_line'], chunk['end_line'], chunk['chunk_type'], chunk_hash))
            
            conn.commit()
            conn.close()
            return True
            
        except Exception as e:
            st.session_state.file_processing_error = f"Error processing {file_path}: {e}"
            return False
    
    def get_relevant_files(self, query: str, max_tokens: int) -> List[Dict]:
        """Get most relevant files within token limit using intelligent selection."""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute("""
            SELECT file_path, extension, tokens, priority, file_type, content_preview
            FROM files 
            WHERE is_processed = TRUE 
            ORDER BY priority DESC, tokens ASC
        """)
        
        files = cursor.fetchall()
        conn.close()
        
        if not files:
            return []
        
        # Simple relevance scoring based on query keywords
        query_keywords = set(query.lower().split())
        selected_files = []
        current_tokens = 0
        
        for file_path, extension, tokens, priority, file_type, content_preview in files:
            if current_tokens + tokens > max_tokens:
                break
                
            filename = os.path.basename(file_path).lower()
            
            # Calculate relevance score
            relevance_score = 0
            if any(keyword in filename for keyword in query_keywords):
                relevance_score += 10
            if extension[1:] in query_keywords:
                relevance_score += 5
            if any(keyword in file_type for keyword in query_keywords):
                relevance_score += 3
            if any(keyword in content_preview.lower() for keyword in query_keywords):
                relevance_score += 7
            
            selected_files.append({
                'file_path': file_path,
                'tokens': tokens,
                'priority': priority,
                'relevance_score': relevance_score,
                'file_type': file_type
            })
            current_tokens += tokens
        
        # Sort by relevance and priority
        selected_files.sort(key=lambda x: (x['relevance_score'], x['priority']), reverse=True)
        return selected_files[:25]
    
    def get_relevant_chunks(self, query: str, file_paths: List[str], max_tokens: int) -> List[Dict]:
        """Get most relevant chunks from selected files."""
        if not st.session_state.embeddings_model:
            return []
        
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # Get chunks from selected files
        placeholders = ','.join(['?' for _ in file_paths])
        cursor.execute(f"""
            SELECT fc.id, fc.content, fc.tokens, fc.start_line, fc.chunk_type, f.file_path
            FROM file_chunks fc
            JOIN files f ON fc.file_id = f.id
            WHERE f.file_path IN ({placeholders})
            ORDER BY fc.tokens ASC
        """, file_paths)
        
        chunks = cursor.fetchall()
        conn.close()
        
        if not chunks:
            return []
        
        # Create query embedding
        query_embedding = st.session_state.embeddings_model.encode([query])[0]
        
        # Calculate similarities for chunks
        relevant_chunks = []
        current_tokens = 0
        
        for chunk_id, content, tokens, start_line, chunk_type, file_path in chunks:
            if current_tokens + tokens > max_tokens:
                break
            
            # Create chunk embedding
            chunk_embedding = st.session_state.embeddings_model.encode([content[:1000]])[0]
            similarity = cosine_similarity([query_embedding], [chunk_embedding])[0][0]
            
            relevant_chunks.append({
                'chunk_id': chunk_id,
                'content': content,
                'tokens': tokens,
                'start_line': start_line,
                'chunk_type': chunk_type,
                'file_path': file_path,
                'similarity': similarity
            })
            current_tokens += tokens
        
        # Sort by similarity
        relevant_chunks.sort(key=lambda x: x['similarity'], reverse=True)
        return relevant_chunks[:15]
    
    def get_codebase_summary(self) -> Dict:
        """Get comprehensive codebase statistics."""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute("SELECT COUNT(*), SUM(lines), SUM(tokens), SUM(size) FROM files WHERE is_processed = TRUE")
        total_files, total_lines, total_tokens, total_size = cursor.fetchone()
        
        cursor.execute("""
            SELECT file_type, COUNT(*) as count, SUM(lines) as lines, SUM(tokens) as tokens
            FROM files 
            WHERE is_processed = TRUE 
            GROUP BY file_type 
            ORDER BY count DESC
        """)
        file_types = {row[0]: {'count': row[1], 'lines': row[2], 'tokens': row[3]} for row in cursor.fetchall()}
        
        cursor.execute("""
            SELECT extension, COUNT(*) as count
            FROM files 
            WHERE is_processed = TRUE 
            GROUP BY extension 
            ORDER BY count DESC
        """)
        extensions = {row[0]: row[1] for row in cursor.fetchall()}
        
        conn.close()
        
        return {
            'total_files': total_files or 0,
            'total_lines': total_lines or 0,
            'total_tokens': total_tokens or 0,
            'total_size': total_size or 0,
            'file_types': file_types,
            'extensions': extensions
        }

@st.cache_resource
def load_embeddings_model():
    """Load sentence transformer model for code embeddings with retry."""
    try:
        model = SentenceTransformer('all-MiniLM-L6-v2')
        return model
    except Exception as e:
        st.error(f"Failed to load embeddings model: {e}")
        # Try again with smaller model
        try:
            model = SentenceTransformer('all-MiniLM-L6-v2')
            return model
        except:
            return None

def estimate_tokens(text: str) -> int:
    """Estimate token count for text with fallback."""
    try:
        encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")
        return len(encoding.encode(str(text)))
    except:
        return len(str(text)) // 4

def get_system_prompt() -> str:
    """Enhanced system prompt for coding tasks."""
    return """You are CodeMaster AI, an expert programming assistant specializing in large codebase analysis and document processing.

**Core Capabilities:**
1. **Code Generation**: Write clean, efficient, well-documented code
2. **Multi-language Support**: Expert in Python, JavaScript, Java, C++, Go, Rust, and more
3. **Large Codebase Analysis**: Understand complex project structures and relationships
4. **Document Processing**: Analyze PDFs, Word docs, PowerPoint, Excel, and other formats
5. **Smart Code Search**: Find relevant code sections using semantic understanding
6. **Optimization**: Suggest performance and architectural improvements
7. **Documentation**: Generate comprehensive explanations and comments

**Response Guidelines:**
- Always specify programming language for code blocks
- Provide context-aware suggestions based on the loaded codebase
- Reference specific files and functions when relevant
- Explain complex logic with clear comments
- Suggest best practices and optimizations
- Keep responses focused and actionable

**Code Quality Standards:**
- Production-ready code with proper error handling
- Follow language-specific conventions and best practices
- Optimize for readability, maintainability, and performance
- Include relevant tests and documentation

Focus on delivering accurate, contextual solutions for large-scale development projects and comprehensive document analysis."""

def manage_context_window(messages: List[Dict], max_tokens: int) -> List[Dict]:
    """Manage conversation context to stay within token limits."""
    if not messages:
        return messages
    
    total_tokens = 0
    kept_messages = []
    
    # Always keep system message
    system_msg = None
    if messages and messages[0].get('role') == 'system':
        system_msg = messages[0]
        total_tokens += estimate_tokens(system_msg['content'])
    
    # Always keep the last user message
    last_user_msg = None
    if messages and messages[-1].get('role') == 'user':
        last_user_msg = messages[-1]
        total_tokens += estimate_tokens(last_user_msg['content'])
    
    # Add messages from most recent backwards
    for msg in reversed(messages[1:-1] if len(messages) > 2 else []):
        msg_tokens = estimate_tokens(msg['content'])
        if total_tokens + msg_tokens <= max_tokens:
            kept_messages.insert(0, msg)
            total_tokens += msg_tokens
        else:
            break
    
    # Reconstruct message list
    final_messages = []
    if system_msg:
        final_messages.append(system_msg)
    final_messages.extend(kept_messages)
    if last_user_msg and last_user_msg not in final_messages:
        final_messages.append(last_user_msg)
    
    return final_messages

def create_intelligent_context(query: str, codebase_handler: EnhancedCodebaseHandler, max_tokens: int) -> str:
    """Create intelligent context for the query within token limits."""
    if not codebase_handler:
        return ""
    
    # Reserve tokens for different parts
    summary_tokens = max_tokens * 0.2
    files_tokens = max_tokens * 0.3
    chunks_tokens = max_tokens * 0.5
    
    context_parts = []
    
    try:
        # 1. Codebase summary
        stats = codebase_handler.get_codebase_summary()
        if stats['total_files'] > 0:
            file_types_str = ', '.join(f"{ftype}({data['count']})" for ftype, data in stats['file_types'].items())
            extensions_str = ', '.join(f"{ext}({count})" for ext, count in list(stats['extensions'].items())[:10])
            
            summary = f"""
**Codebase Overview:**
- Total Files: {stats['total_files']:,}
- Total Lines: {stats['total_lines']:,}
- Total Tokens: {stats['total_tokens']:,}
- File Types: {file_types_str}
- Extensions: {extensions_str}
"""
            if estimate_tokens(summary) <= summary_tokens:
                context_parts.append(summary)
        
        # 2. Relevant files
        relevant_files = codebase_handler.get_relevant_files(query, int(files_tokens))
        if relevant_files:
            files_context = "\n**Relevant Files:**\n"
            for file_info in relevant_files[:8]:  # Reduced to top 8 files
                files_context += f"- {file_info['file_path']} ({file_info['file_type']}, {file_info['tokens']:,} tokens)\n"
            
            if estimate_tokens(files_context) <= files_tokens:
                context_parts.append(files_context)
        
        # 3. Relevant code chunks (only if we have embeddings model)
        if st.session_state.embeddings_model and relevant_files:
            file_paths = [f['file_path'] for f in relevant_files[:4]]  # Top 4 files
            relevant_chunks = codebase_handler.get_relevant_chunks(query, file_paths, int(chunks_tokens))
            
            if relevant_chunks:
                chunks_context = "\n**Most Relevant Code/Content Sections:**\n"
                for i, chunk in enumerate(relevant_chunks[:4]):  # Top 4 chunks
                    chunks_context += f"\n**File: {chunk['file_path']} (Line {chunk['start_line']})**\n"
                    content = chunk['content']
                    if len(content) > 800:  # More aggressive truncation
                        content = content[:400] + "\n... (truncated) ...\n" + content[-400:]
                    chunks_context += f"```\n{content}\n```\n"
                
                chunks_tokens_actual = estimate_tokens(chunks_context)
                if chunks_tokens_actual <= chunks_tokens:
                    context_parts.append(chunks_context)
    
    except Exception as e:
        st.error(f"Error creating context: {str(e)}")
    
    return "\n".join(context_parts)

def call_llm_api(messages: List[Dict], api_key: str, model: str, temperature: float, max_tokens: int) -> Optional[str]:
    """Call Groq API with enhanced error handling and retry logic."""
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "User-Agent": "LLM-Code-Assistant/1.0"
    }
    
    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_tokens,
        "stream": False
    }
    
    max_retries = 3
    backoff_factor = 1
    timeout = 30  # seconds
    
    for attempt in range(max_retries):
        try:
            response = requests.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers=headers,
                json=payload,
                timeout=timeout
            )
            
            response.raise_for_status()  # Raise exception for HTTP errors
            
            if response.status_code == 200:
                result = response.json()
                return result['choices'][0]['message']['content']
            elif response.status_code == 429:
                wait_time = backoff_factor * (2 ** attempt)
                st.warning(f"Rate limit exceeded. Retrying in {wait_time} seconds...")
                time.sleep(wait_time)
                continue
            else:
                error_msg = f"API Error {response.status_code}: {response.text}"
                st.error(error_msg)
                return None
                
        except requests.exceptions.Timeout:
            st.warning(f"Request timed out. Attempt {attempt+1}/{max_retries}")
            if attempt < max_retries - 1:
                time.sleep(backoff_factor * (2 ** attempt))
            else:
                st.error("API request timed out after multiple attempts")
                return None
        except requests.exceptions.RequestException as e:
            st.error(f"Request failed: {str(e)}")
            return None
        except Exception as e:
            st.error(f"Unexpected error: {str(e)}")
            return None
    
    st.error("Failed to get response after multiple attempts")
    return None

def run_code(code: str, language: str) -> Dict:
    """Execute code in various programming languages."""
    if language not in SUPPORTED_LANGUAGES:
        return {"error": f"Language {language} not supported"}
    
    lang_config = SUPPORTED_LANGUAGES[language]
    extension = lang_config['extension']
    runner = lang_config['runner']
    
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            # Create temporary file
            file_path = os.path.join(temp_dir, f"temp_code{extension}")
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(code)
            
            # Execute based on language
            if language == 'python':
                result = subprocess.run([sys.executable, file_path], 
                                      capture_output=True, text=True, timeout=30)
            elif language == 'javascript':
                result = subprocess.run(['node', file_path], 
                                      capture_output=True, text=True, timeout=30)
            elif language == 'java':
                # Compile first
                compile_result = subprocess.run(['javac', file_path], 
                                              capture_output=True, text=True, timeout=30)
                if compile_result.returncode != 0:
                    return {"error": f"Compilation error: {compile_result.stderr}"}
                
                # Extract class name and run
                class_name = os.path.splitext(os.path.basename(file_path))[0]
                result = subprocess.run(['java', '-cp', temp_dir, class_name], 
                                      capture_output=True, text=True, timeout=30)
            elif language in ['cpp', 'c']:
                # Compile first
                executable_path = os.path.join(temp_dir, 'temp_executable')
                compiler = 'g++' if language == 'cpp' else 'gcc'
                compile_result = subprocess.run([compiler, file_path, '-o', executable_path], 
                                              capture_output=True, text=True, timeout=30)
                if compile_result.returncode != 0:
                    return {"error": f"Compilation error: {compile_result.stderr}"}
                
                result = subprocess.run([executable_path], 
                                      capture_output=True, text=True, timeout=30)
            else:
                result = subprocess.run([runner, file_path], 
                                      capture_output=True, text=True, timeout=30)
            
            return {
                "stdout": result.stdout,
                "stderr": result.stderr,
                "returncode": result.returncode
            }
            
    except subprocess.TimeoutExpired:
        return {"error": "Code execution timed out (30s limit)"}
    except FileNotFoundError:
        return {"error": f"Language runtime not found: {runner}"}
    except Exception as e:
        return {"error": f"Execution error: {str(e)}"}

def process_uploaded_files(uploaded_files, codebase_handler: EnhancedCodebaseHandler) -> bool:
    """Process multiple uploaded files with progress tracking and error handling."""
    if not uploaded_files:
        return False
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    st.session_state.file_processing_error = None
    
    processed_count = 0
    total_files = len(uploaded_files)
    
    with tempfile.TemporaryDirectory() as temp_dir:
        for i, uploaded_file in enumerate(uploaded_files):
            try:
                # Save uploaded file
                file_path = os.path.join(temp_dir, uploaded_file.name)
                with open(file_path, "wb") as f:
                    f.write(uploaded_file.getbuffer())
                
                # Process file
                status_text.text(f"Processing: {uploaded_file.name}")
                if codebase_handler.process_file(file_path):
                    processed_count += 1
                else:
                    if st.session_state.file_processing_error:
                        st.warning(st.session_state.file_processing_error)
                
                progress_bar.progress((i + 1) / total_files)
                
            except Exception as e:
                st.error(f"Error processing {uploaded_file.name}: {e}")
    
    status_text.text(f"Processed {processed_count}/{total_files} files successfully")
    progress_bar.empty()
    status_text.empty()
    
    if processed_count > 0:
        # Load embeddings after processing
        st.session_state.embeddings_model = load_embeddings_model()
    
    return processed_count > 0

def process_directory(directory_path: str, codebase_handler: EnhancedCodebaseHandler) -> bool:
    """Process all files in a directory recursively."""
    if not os.path.exists(directory_path):
        st.error(f"Directory not found: {directory_path}")
        return False
    
    # Find all processable files
    files_to_process = []
    for root, dirs, files in os.walk(directory_path):
        # Skip common directories
        dirs[:] = [d for d in dirs if not d.startswith(('.git', '__pycache__', 'node_modules', '.venv', 'venv', '.vs', '.vscode'))]
        
        for file in files:
            file_path = os.path.join(root, file)
            if Path(file_path).suffix.lower() in ALL_EXTENSIONS:
                files_to_process.append(file_path)
    
    if not files_to_process:
        st.warning("No processable files found in directory")
        return False
    
    # Process files with progress tracking
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    processed_count = 0
    total_files = len(files_to_process)
    
    for i, file_path in enumerate(files_to_process):
        try:
            status_text.text(f"Processing: {os.path.basename(file_path)} ({i+1}/{total_files})")
            if codebase_handler.process_file(file_path):
                processed_count += 1
            else:
                if st.session_state.file_processing_error:
                    st.warning(st.session_state.file_processing_error)
            
            progress_bar.progress((i + 1) / total_files)
            
        except Exception as e:
            st.error(f"Error processing {file_path}: {e}")
    
    status_text.text(f"Processed {processed_count}/{total_files} files successfully")
    progress_bar.empty()
    status_text.empty()
    
    if processed_count > 0:
        # Load embeddings after processing
        st.session_state.embeddings_model = load_embeddings_model()
    
    return processed_count > 0

def main():
    """Main Streamlit application with fixes for API issues."""
    
    # Initialize enhanced codebase handler
    if st.session_state.codebase_handler is None:
        st.session_state.codebase_handler = EnhancedCodebaseHandler()
    
    # Load embeddings model on startup
    if st.session_state.embeddings_model is None:
        st.session_state.embeddings_model = load_embeddings_model()
    
    # Original header
    st.markdown('<div class="main-header"><h1>LLM Code Assistant</h1><p>Advanced code analysis and general tasks</p></div>', unsafe_allow_html=True)
    
    # Sidebar
    with st.sidebar:
        st.header("⚙️ Configuration")
        
        # API Settings
        api_key = st.text_input("Groq API Key", type="password", value=st.session_state.api_key)
        if api_key != st.session_state.api_key:
            st.session_state.api_key = api_key
        
        # Model Selection
        model_options = [f"{info['name']}" for info in GROQ_MODELS.values()]
        current_model_name = GROQ_MODELS[st.session_state.model]["name"]
        
        selected_model = st.selectbox(
            "🤖 Model",
            model_options,
            index=model_options.index(current_model_name)
        )
        
        # Update session state model
        for key, value in GROQ_MODELS.items():
            if value["name"] == selected_model:
                st.session_state.model = key
                st.session_state.max_context_tokens = min(value["context"] - 4000, 120000)
                break
        
        st.info(f"Context: {GROQ_MODELS[st.session_state.model]['context']:,} tokens")
        
        # Generation Parameters
        st.subheader("🎛️ Parameters")
        st.session_state.temperature = st.slider("Temperature", 0.0, 1.0, st.session_state.temperature, 0.1)
        st.session_state.max_tokens = st.slider("Max Tokens", 100, 8000, st.session_state.max_tokens, 100)
        
        # Enhanced Codebase Management
        st.header("Codebase Management")
        
        # File Upload
        st.subheader("Upload Files")
        uploaded_files = st.file_uploader(
            "Choose files (supports code, documents, PDFs, etc.)",
            accept_multiple_files=True,
            type=list(set(ext[1:] for ext in ALL_EXTENSIONS))  # All supported extensions
        )
        
        if uploaded_files and st.button("Upload & Process"):
            with st.spinner("Processing uploaded files..."):
                if process_uploaded_files(uploaded_files, st.session_state.codebase_handler):
                    st.success(f"Successfully processed {len(uploaded_files)} files!")
                    st.rerun()
        
        # Directory Processing
        st.subheader("Process Directory")
        directory_path = st.text_input("📂 Directory Path", placeholder="/path/to/your/code")
        if directory_path and st.button("🔄 Process Directory"):
            with st.spinner("Processing directory..."):
                if process_directory(directory_path, st.session_state.codebase_handler):
                    st.success("Directory processed successfully!")
                    st.rerun()
        
        # Codebase Statistics
        stats = st.session_state.codebase_handler.get_codebase_summary()
        if stats['total_files'] > 0:
            st.subheader("📊 Codebase Stats")
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Files", f"{stats['total_files']:,}")
                st.metric("Tokens", f"{stats['total_tokens']:,}")
            with col2:
                st.metric("Lines", f"{stats['total_lines']:,}")
                st.metric("Size", f"{stats['total_size']/1024/1024:.1f}MB")
            
            # File types breakdown
            if stats['file_types']:
                st.write("**File Types:**")
                for ftype, data in stats['file_types'].items():
                    st.write(f"• {ftype}: {data['count']} files")
        
        # Clear Codebase
        if st.button("🗑️ Clear Codebase", type="secondary"):
            if os.path.exists(st.session_state.codebase_handler.db_path):
                os.remove(st.session_state.codebase_handler.db_path)
            st.session_state.codebase_handler = EnhancedCodebaseHandler()
            st.session_state.messages = []
            st.session_state.embeddings_model = None
            st.success("Codebase cleared!")
            st.rerun()
        
        # Add Core Capabilities section at the bottom
        st.markdown("---")
        st.markdown("""
        <div class="core-capabilities">
            <h4>Core Capabilities:</h4>
            <div class="capability-item">
                <span class="capability-icon"></span>
                <span>Multi-language Code Generation</span>
            </div>
            <div class="capability-item">
                <span class="capability-icon"></span>
                <span>Smart Codebase Analysis</span>
            </div>
            <div class="capability-item">
                <span class="capability-icon"></span>
                <span>Document Processing (PDF, Word, Excel)</span>
            </div>
            <div class="capability-item">
                <span class="capability-icon"></span>
                <span>Semantic Code Search</span>
            </div>
            <div class="capability-item">
                <span class="capability-icon"></span>
                <span>Performance Optimization</span>
            </div>
            <div class="capability-item">
                <span class="capability-icon"></span>
                <span>Code Execution & Testing</span>
            </div>
            <div class="capability-item">
                <span class="capability-icon"></span>
                <span>Comprehensive Documentation</span>
            </div>
        </div>
        """, unsafe_allow_html=True)
    
    # Main Chat Interface
    st.header("Chat with your LLM")
    
    # Display chat messages
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            # Handle code blocks
            content = message["content"]
            if "```" in content:
                parts = content.split("```")
                for i, part in enumerate(parts):
                    if i % 2 == 0:
                        if part.strip():
                            st.markdown(part)
                    else:
                        # This is a code block
                        lines = part.split('\n')
                        language = lines[0].strip() if lines else ''
                        code_content = '\n'.join(lines[1:]) if len(lines) > 1 else part
                        
                        # Display language and run button
                        if language in SUPPORTED_LANGUAGES:
                            col1, col2 = st.columns([3, 1])
                            with col1:
                                st.code(code_content, language=language)
                            with col2:
                                if st.button(f"▶️ Run", key=f"run_{hash(code_content)}"):
                                    result = run_code(code_content, language)
                                    if "error" in result:
                                        st.error(result["error"])
                                    else:
                                        if result["stdout"]:
                                            st.success("Output:")
                                            st.text(result["stdout"])
                                        if result["stderr"]:
                                            st.warning("Errors:")
                                            st.text(result["stderr"])
                        else:
                            st.code(code_content, language=language if language else None)
            else:
                st.markdown(content)
    
    # Chat input
    if prompt := st.chat_input("Ask about your codebase or request code assistance..."):
        if not st.session_state.api_key:
            st.error("Please enter your Groq API key in the sidebar.")
            st.stop()
        
        # Add user message
        st.session_state.messages.append({"role": "user", "content": prompt})
        
        # Display user message
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Generate assistant response
        with st.chat_message("assistant"):
            with st.spinner("Analyzing..."):
                try:
                    # Create intelligent context (with reduced size)
                    context = create_intelligent_context(
                        prompt, 
                        st.session_state.codebase_handler, 
                        min(8000, st.session_state.max_context_tokens // 3)  # Reduced context size
                    )
                    
                    # Prepare messages
                    system_prompt = get_system_prompt()
                    if context:
                        system_prompt += f"\n\n**Current Codebase Context:**\n{context}"
                    
                    api_messages = [{"role": "system", "content": system_prompt}]
                    api_messages.extend(st.session_state.messages[-6:])  # Only last 6 messages
                    
                    # Estimate tokens
                    total_tokens = sum(estimate_tokens(msg['content']) for msg in api_messages)
                    st.caption(f"Estimated tokens: {total_tokens}/{st.session_state.max_context_tokens}")
                    
                    # Call API
                    response = call_llm_api(
                        api_messages,
                        st.session_state.api_key,
                        st.session_state.model,
                        st.session_state.temperature,
                        st.session_state.max_tokens
                    )
                    
                    if response:
                        # Display response
                        st.session_state.messages.append({"role": "assistant", "content": response})
                        st.markdown(response)
                    else:
                        st.error("Failed to get response from the API. Please check your API key and network connection.")
                        st.session_state.messages.append({"role": "assistant", "content": "I'm having trouble connecting to the API. Please check your settings."})
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
                    st.session_state.messages.append({"role": "assistant", "content": "An error occurred while processing your request."})

if __name__ == "__main__":
    main()
